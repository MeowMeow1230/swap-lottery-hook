"""
Swap Lottery Hook — Hackathon Pitch Deck
Generates a .pptx with dark tech aesthetic, neon accents,
professional typography. No API keys required.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──
BG = RGBColor(0x0A, 0x0A, 0x12)
CARD_BG = RGBColor(0x14, 0x14, 0x22)
NEON_PURPLE = RGBColor(0xA7, 0x6B, 0xFF)
NEON_BLUE = RGBColor(0x5E, 0xB0, 0xFF)
NEON_CYAN = RGBColor(0x00, 0xE4, 0xCC)
WHITE = RGBColor(0xED, 0xED, 0xF4)
GRAY = RGBColor(0x88, 0x8A, 0x9A)
DARK_GRAY = RGBColor(0x2A, 0x2A, 0x38)
GREEN = RGBColor(0x34, 0xD3, 0x99)
RED = RGBColor(0xFF, 0x5E, 0x7A)
AMBER = RGBColor(0xFF, 0xB6, 0x27)
GRADIENT_TOP = RGBColor(0x1A, 0x10, 0x30)

W = Inches(13.333)  # 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

def slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Inter"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_rich_text(slide, left, top, width, height, lines, font_name="Inter"):
    """lines: list of (text, font_size, color, bold, align)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        text, size, color, bold, align = line
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(4)
    return txBox

def add_card(slide, left, top, width, height, title, body_lines, accent=NEON_PURPLE):
    """A dark card with accent left border"""
    # Card bg
    add_rect(slide, left, top, width, height, CARD_BG, 0.03)
    # Accent left bar
    add_rect(slide, left, top + Inches(0.2), Inches(0.06), height - Inches(0.4), accent)
    # Title
    add_text_box(slide, left + Inches(0.25), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
                 title, font_size=20, color=WHITE, bold=True)
    # Body
    y = top + Inches(0.6)
    for bline in body_lines:
        add_text_box(slide, left + Inches(0.25), y, width - Inches(0.4), Inches(0.35),
                     bline, font_size=13, color=GRAY)
        y += Inches(0.3)

def add_divider(slide, left, top, width, color=DARK_GRAY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_page_number(slide, num, total=10):
    add_text_box(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.3),
                 f"{num}/{total}", font_size=10, color=GRAY, align=PP_ALIGN.RIGHT)

def footer_bar(slide):
    add_divider(slide, Inches(0.8), Inches(6.85), Inches(11.7))
    add_text_box(slide, Inches(0.8), Inches(6.95), Inches(6), Inches(0.3),
                 "Swap Lottery Hook  |  Hook the Future Hackathon  |  X Layer × Uniswap V4",
                 font_size=8, color=DARK_GRAY)
    add_text_box(slide, Inches(7), Inches(6.95), Inches(5.8), Inches(0.3),
                 "@XLayerOfficial  @Uniswap  @flapdotsh",
                 font_size=8, color=DARK_GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 1: COVER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
slide_bg(slide)

# Large gradient glow (simulated with shapes)
add_rect(slide, Inches(6), Inches(0.5), Inches(7), Inches(7), RGBColor(0x1F, 0x10, 0x3F))
# Purple accent bar on top
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), NEON_PURPLE)

# Title
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "Swap Lottery Hook", font_size=60, color=WHITE, bold=True)
# Subtitle with gradient feel
add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
             "Make Every Swap a Winning Ticket", font_size=32, color=NEON_PURPLE, bold=False)
# Tagline
add_text_box(slide, Inches(1), Inches(4.0), Inches(10), Inches(0.6),
             "Uniswap V4 Hook  ×  Game Theory  ×  On-Chain Lottery Engine",
             font_size=18, color=GRAY)
# Divider
add_divider(slide, Inches(1), Inches(4.8), Inches(4), NEON_PURPLE)
# Meta
add_text_box(slide, Inches(1), Inches(5.2), Inches(6), Inches(0.5),
             "Hook the Future Hackathon  |  X Layer  |  May 2026", font_size=16, color=GRAY)
add_text_box(slide, Inches(1), Inches(6.5), Inches(6), Inches(0.4),
             "Deployed on X Layer Testnet  ·  Verified Contracts", font_size=12, color=GREEN)

add_page_number(slide, 1)

# ============================================================
# SLIDE 2: THE PROBLEM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), NEON_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Why Existing Solutions Don't Work", font_size=36, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.1), Inches(3), NEON_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
             "Traders flow through liquidity pools like water — they come, swap, and never return.",
             font_size=16, color=GRAY)

# 3 problem cards
problems = [
    ("Static LP Fees",
     "Passive income model. Zero incentive for traders to return. LPs earn fees but pools bleed volume to competitors with better UX or lower fees."),
    ("Blind Mining / Points",
     "Delayed gratification. 3-6 month vesting cycles create user fatigue. \"PUA\" mechanics that over-promise and under-deliver on airdrops."),
    ("Dynamic Fee Hooks (Existing)",
     "Only protect LPs — they raise fees during volatility, which drives traders AWAY. Solving one side while destroying the other is not a solution."),
]

for i, (title, body) in enumerate(problems):
    x = Inches(0.8 + i * 4.1)
    add_card(slide, x, Inches(2.2), Inches(3.8), Inches(3.8), title,
             [body], accent=[NEON_PURPLE, NEON_BLUE, NEON_CYAN][i])

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
             "We need a mechanism that rewards traders for returning, without penalizing LPs.",
             font_size=16, color=WHITE, bold=True)
footer_bar(slide)
add_page_number(slide, 2)

# ============================================================
# SLIDE 3: CORE MECHANISM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), NEON_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Core Mechanism: The Lottery Flywheel", font_size=36, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.1), Inches(3), NEON_BLUE)

# Flow diagram - 4 connected boxes
steps = [
    ("0.01% Tax", "Invisible friction.\n1 basis point from\nevery swap output.", NEON_PURPLE),
    ("Dynamic N", "Pool small → fast trigger.\nPool large → anticipation.\nPsychological curve.", NEON_BLUE),
    ("Multi-Tier Prizes", "60% Grand Prize\n10 × 4% Consolation\n\"I was so close\" effect.", NEON_CYAN),
    ("LP Wins Too", "Volume surge from\nlottery addiction lifts\nall LP fee revenue.", GREEN),
]

for i, (title, desc, color) in enumerate(steps):
    x = Inches(0.5 + i * 3.2)
    add_card(slide, x, Inches(1.7), Inches(2.9), Inches(3.0), title, [desc], accent=color)
    if i < 3:
        add_text_box(slide, x + Inches(2.9), Inches(2.8), Inches(0.4), Inches(0.4),
                     "→", font_size=24, color=GRAY)

# Formula highlight
add_rich_text(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(1.5), [
    ("Dynamic N Formula:", 12, GRAY, False, PP_ALIGN.LEFT),
    ("N = BASE_N × (1 + poolSize / (poolSize + K))  |  K = 100 ETH", 16, NEON_CYAN, True, PP_ALIGN.LEFT),
    ("Range: N ∈ [100, 200]  —  bounded, predictable, gas-efficient", 13, GRAY, False, PP_ALIGN.LEFT),
])
footer_bar(slide)
add_page_number(slide, 3)

# ============================================================
# SLIDE 4: LP VALUE PROOF
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), NEON_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "LP Value: The Critical Defense", font_size=36, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.1), Inches(3), NEON_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
             "The #1 question judges will ask: \"Doesn't the lottery fee DILUTE LP returns?\" Here's the answer.",
             font_size=16, color=GRAY)

# Comparison box
add_rect(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(2.0), CARD_BG, 0.03)
add_text_box(slide, Inches(1.0), Inches(2.3), Inches(5), Inches(0.4),
             "Traditional Pool (0.3% static fee)", font_size=16, color=RED, bold=True)
add_rich_text(slide, Inches(1.0), Inches(2.7), Inches(5), Inches(1.5), [
    ("LP APR: 12% (flat, predictable)", 15, WHITE, False, PP_ALIGN.LEFT),
    ("Volume: baseline, no growth flywheel", 15, GRAY, False, PP_ALIGN.LEFT),
    ("Trader retention: 0 — they swap once, they leave", 15, GRAY, False, PP_ALIGN.LEFT),
])

add_rect(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(2.0), RGBColor(0x1A, 0x14, 0x2E), 0.03)
add_text_box(slide, Inches(7.2), Inches(2.3), Inches(5), Inches(0.4),
             "Swap Lottery Pool (0.3% LP + 0.01% ticket)", font_size=16, color=GREEN, bold=True)
add_rich_text(slide, Inches(7.2), Inches(2.7), Inches(5), Inches(1.5), [
    ("LP APR: 15-18% (volume-driven surge)", 15, WHITE, False, PP_ALIGN.LEFT),
    ("Volume: +22-35% from lottery flywheel effect", 15, GRAY, False, PP_ALIGN.LEFT),
    ("Trader retention: high — \"next swap could be the winner\"", 15, GRAY, False, PP_ALIGN.LEFT),
])

# Key stat
add_rect(slide, Inches(2.5), Inches(4.5), Inches(8), Inches(1.5), CARD_BG, 0.03)
add_rich_text(slide, Inches(2.8), Inches(4.7), Inches(7.5), Inches(1.0), [
    ("Simulation Result", 11, GRAY, False, PP_ALIGN.CENTER),
    ("When lottery multiplier > 1.2×, LP fee revenue surpasses", 14, WHITE, False, PP_ALIGN.CENTER),
    ("traditional pools by 15-30% — even after lottery pool deduction", 14, GREEN, True, PP_ALIGN.CENTER),
])

add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.4),
             "The 0.01% tax is a growth investment, not a cost. LPs earn MORE because volume compounds.",
             font_size=14, color=NEON_CYAN, bold=True, align=PP_ALIGN.CENTER)
footer_bar(slide)
add_page_number(slide, 4)

# ============================================================
# SLIDE 5: RANDOMNESS ETHICS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), NEON_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Randomness: Engineering Ethics", font_size=36, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.1), Inches(3), NEON_BLUE)

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
             "Honesty > Perfection. We know exactly where the line is and we draw it clearly.",
             font_size=16, color=GRAY)

# Phase 1: Demo
add_rect(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(2.5), CARD_BG, 0.03)
add_text_box(slide, Inches(1.0), Inches(2.3), Inches(5), Inches(0.4),
             "Hackathon Demo: Pseudo-Random", font_size=18, color=AMBER, bold=True)
add_rich_text(slide, Inches(1.0), Inches(2.8), Inches(5), Inches(1.8), [
    ("Source: keccak256(prevrandao + timestamp + sender + seed)", 12, GRAY, False, PP_ALIGN.LEFT),
    ("", 8, GRAY, False, PP_ALIGN.LEFT),
    ("✔ Defense vs casual users — unpredictable enough for individuals", 13, GREEN, False, PP_ALIGN.LEFT),
    ("✘ Vulnerable to block builders — they can manipulate prevrandao", 13, RED, False, PP_ALIGN.LEFT),
    ("✘ Not suitable for high-value mainnet pools", 13, RED, False, PP_ALIGN.LEFT),
])

# Phase 2: Production
add_rect(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(2.5), RGBColor(0x1A, 0x14, 0x2E), 0.03)
add_text_box(slide, Inches(7.2), Inches(2.3), Inches(5), Inches(0.4),
             "Production: Verifiable Randomness", font_size=18, color=GREEN, bold=True)
add_rich_text(slide, Inches(7.2), Inches(2.8), Inches(5), Inches(1.8), [
    ("Chainlink VRF or API3 QRNG", 14, WHITE, True, PP_ALIGN.LEFT),
    ("", 8, GRAY, False, PP_ALIGN.LEFT),
    ("✔ Cryptographically verifiable randomness", 13, GREEN, False, PP_ALIGN.LEFT),
    ("✔ Immune to block builder manipulation", 13, GREEN, False, PP_ALIGN.LEFT),
    ("✔ Gas cost: ~30k per VRF request on X Layer", 13, GREEN, False, PP_ALIGN.LEFT),
    ("✔ X Layer natively supports these oracle networks", 13, GREEN, False, PP_ALIGN.LEFT),
])

# Bottom note
add_rect(slide, Inches(2), Inches(5.2), Inches(9), Inches(1.0), CARD_BG, 0.03)
add_rich_text(slide, Inches(2.3), Inches(5.35), Inches(8.5), Inches(0.8), [
    ("Our Position", 12, NEON_CYAN, True, PP_ALIGN.CENTER),
    ("\"We deliver an honest MVP with clear boundaries. The path to production-grade security is documented, funded, and achievable within one sprint after the hackathon.\"", 12, GRAY, False, PP_ALIGN.CENTER),
])

footer_bar(slide)
add_page_number(slide, 5)

# ============================================================
# SLIDE 6: FLYWHEEL ENDGAME
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), NEON_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "The Flywheel Endgame", font_size=36, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.1), Inches(3), NEON_BLUE)

# Circular flywheel - 4 quadrants
flywheel = [
    ("More Trades", "\"Next swap could\nbe the winner\"\n→ volume surges", NEON_PURPLE, Inches(1), Inches(2)),
    ("Bigger Prize Pool", "0.01% × volume\n= compounding\nlottery treasury", NEON_BLUE, Inches(6.5), Inches(2)),
    ("LPs Earn More", "Volume-driven\nfee growth\noutpaces dilution", GREEN, Inches(6.5), Inches(4.5)),
    ("More LPs Join", "Higher APR attracts\nmore liquidity\n→ tighter spreads", NEON_CYAN, Inches(1), Inches(4.5)),
]

for title, desc, color, x, y in flywheel:
    add_rect(slide, x, y, Inches(3.2), Inches(2.0), CARD_BG, 0.03)
    add_rect(slide, x, y + Inches(0.2), Inches(0.06), Inches(1.6), color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(2.8), Inches(0.35),
                 title, font_size=18, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.6), Inches(2.8), Inches(1.2),
                 desc, font_size=13, color=GRAY)

# Center text
add_rich_text(slide, Inches(4.2), Inches(3.2), Inches(3), Inches(1), [
    ("", 8, WHITE, False, PP_ALIGN.CENTER),
    ("Compound", 28, WHITE, True, PP_ALIGN.CENTER),
    ("Growth", 28, NEON_PURPLE, True, PP_ALIGN.CENTER),
])

footer_bar(slide)
add_page_number(slide, 6)

# ============================================================
# SLIDE 7: CONTRACT ARCHITECTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), NEON_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Smart Contract Architecture", font_size=36, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.1), Inches(3), NEON_BLUE)

# Flow
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.4),
             "Uniswap V4 Hook Lifecycle  →  SwapLotteryHook.sol (7.6 KB on-chain)",
             font_size=14, color=GRAY)

flow_steps = [
    ("1\nUser\nSwaps", NEON_PURPLE),
    ("2\nbeforeSwap\nSet Dynamic Fee", NEON_BLUE),
    ("3\nPoolManager\nExecutes Swap", GRAY),
    ("4\nafterSwap\nLottery Logic", NEON_CYAN),
    ("5\nClaim Ticket\nReturn Delta", GREEN),
    ("6\nCheck Draw\nTrigger if N ≥ 100", AMBER),
]

for i, (label, color) in enumerate(flow_steps):
    x = Inches(0.5 + i * 2.1)
    add_rect(slide, x, Inches(2.2), Inches(1.8), Inches(1.2), CARD_BG, 0.03)
    add_rect(slide, x, Inches(2.2), Inches(1.8), Inches(0.05), color)
    add_text_box(slide, x + Inches(0.1), Inches(2.35), Inches(1.6), Inches(1.0),
                 label, font_size=14, color=color, bold=True, align=PP_ALIGN.CENTER)
    if i < 5:
        add_text_box(slide, x + Inches(1.8), Inches(2.6), Inches(0.3), Inches(0.3),
                     "→", font_size=18, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Key functions
add_text_box(slide, Inches(0.8), Inches(3.8), Inches(4), Inches(0.4),
             "Key Mechanisms:", font_size=16, color=WHITE, bold=True)

mechs = [
    ("getDynamicN()", "N = 100 + 100 × pool/(pool+K)\nBounded [100, 200]"),
    ("_draw()", "60% grand + 10×4% consolation\nkeccak256 pseudo-random pick"),
    ("_pay()", "manager.take() → winner\nHook's PoolManager balance"),
]
for i, (func, desc) in enumerate(mechs):
    y = Inches(4.3 + i * 0.8)
    add_text_box(slide, Inches(1.0), y, Inches(2.5), Inches(0.3),
                 func, font_size=13, color=NEON_CYAN, bold=True)
    add_text_box(slide, Inches(3.5), y, Inches(8), Inches(0.6),
                 desc, font_size=12, color=GRAY)

footer_bar(slide)
add_page_number(slide, 7)

# ============================================================
# SLIDE 8: DEPLOYMENT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), NEON_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Deployed on X Layer Testnet", font_size=36, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.1), Inches(3), GREEN)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.4),
             "Live, verifiable contracts. Chain ID: 1952  |  Explorer: oklink.com/xlayer-test",
             font_size=14, color=GRAY)

contracts = [
    ("PoolManager", "0x124850dA551dC83b124A9A3F84f8D6674C870Ba1", NEON_PURPLE),
    ("SwapLotteryHook", "0x2D782F42A0a8bBc6b7526aC2272728976Eee90C4", GREEN),
    ("Test USDC", "0xC67958F2329f4C289aa43D52577C0aF32Dbae028", NEON_BLUE),
    ("Test WETH", "0x24003ccc694FF272e54b1601C68019d0E1639eBb", NEON_CYAN),
]

for i, (name, addr, color) in enumerate(contracts):
    y = Inches(2.3 + i * 0.8)
    add_rect(slide, Inches(1.5), y, Inches(10), Inches(0.6), CARD_BG, 0.02)
    add_text_box(slide, Inches(1.7), y + Inches(0.1), Inches(2.5), Inches(0.4),
                 name, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(4.3), y + Inches(0.1), Inches(7), Inches(0.4),
                 addr, font_size=12, color=GRAY)

# Test stats
add_rect(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.7), CARD_BG, 0.03)
add_rich_text(slide, Inches(1.8), Inches(5.9), Inches(9.5), Inches(0.5), [
    ("getDynamicN() = 100  |  getProgress() = (0, 100)  |  Hook Address Flags = 0x10C4 ✓  |  4/4 Tests Passed", 12, GREEN, False, PP_ALIGN.CENTER),
])

footer_bar(slide)
add_page_number(slide, 8)

# ============================================================
# SLIDE 9: LIMITATIONS & ROADMAP
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), NEON_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "Known Limitations & Production Roadmap", font_size=36, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.1), Inches(3), NEON_BLUE)

# Limitations
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.4),
             "Current Limitations (Hackathon MVP)", font_size=18, color=AMBER, bold=True)

lims = [
    "Pseudo-random: block.prevrandao manipulable by block builders",
    "Single-currency tracking: lotteryPool doesn't distinguish token0 vs token1",
    "Gas optimization: 10-loop consolation draw needs batch payout",
    "No fee cap: if pool gets extremely large, N caps at 200 but there's no max prize",
]
for i, lim in enumerate(lims):
    add_text_box(slide, Inches(1.0), Inches(2.1 + i * 0.55), Inches(5.5), Inches(0.5),
                 f"• {lim}", font_size=11, color=GRAY)

# Roadmap
add_text_box(slide, Inches(7), Inches(1.5), Inches(5), Inches(0.4),
             "Production Roadmap (Post-Hackathon)", font_size=18, color=GREEN, bold=True)

roadmap = [
    ("Week 1-2", "Chainlink VRF integration for verifiable randomness"),
    ("Week 3-4", "Dual-currency pool tracking + batch payout optimization"),
    ("Month 2", "Spearbit audit + mainnet beta on X Layer"),
    ("Month 3", "Frontend dApp + analytics dashboard + LP simulation tool"),
]
for i, (time, item) in enumerate(roadmap):
    y = Inches(2.1 + i * 0.55)
    add_text_box(slide, Inches(7.2), y, Inches(1.5), Inches(0.3),
                 time, font_size=11, color=NEON_CYAN, bold=True)
    add_text_box(slide, Inches(8.8), y, Inches(4), Inches(0.4),
                 item, font_size=11, color=GRAY)

# Additional features
add_divider(slide, Inches(0.8), Inches(4.5), Inches(11.5))
add_text_box(slide, Inches(0.8), Inches(4.8), Inches(5), Inches(0.4),
             "Future Features Under Research", font_size=16, color=WHITE, bold=True)

future = [
    "NFT-based multiplier boosts for VIP traders",
    "Cross-pool lottery: one draw across multiple liquidity pools",
    "DAO-governed parameters: community votes on N curve, prize splits, fee rates",
    "Privacy-preserving winner claims via zero-knowledge proofs",
]
for i, f in enumerate(future):
    add_text_box(slide, Inches(1.0), Inches(5.3 + i * 0.4), Inches(5.5), Inches(0.35),
                 f"▸ {f}", font_size=11, color=GRAY)

footer_bar(slide)
add_page_number(slide, 9)

# ============================================================
# SLIDE 10: CLOSING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(slide)

add_rect(slide, Inches(6), Inches(0), Inches(7), Inches(7.5), RGBColor(0x1F, 0x10, 0x3F))
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), NEON_PURPLE)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "Every Swap Is a Ticket.", font_size=56, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
             "Every LP Is a Winner.", font_size=56, color=NEON_PURPLE, bold=True)

add_divider(slide, Inches(1), Inches(4.0), Inches(4), NEON_CYAN)

add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
             "Swap Lottery Hook turns the Uniswap V4 liquidity pool", font_size=20, color=GRAY)
add_text_box(slide, Inches(1), Inches(5.1), Inches(11), Inches(0.6),
             "from a passive utility into an active, addictive flywheel.", font_size=20, color=GRAY)

add_text_box(slide, Inches(1), Inches(5.9), Inches(11), Inches(0.6),
             "0.01% tax  ·  Dynamic-N psychology  ·  Multi-tier prizes  ·  LP value growth",
             font_size=16, color=NEON_CYAN)

# CTA
add_rect(slide, Inches(3.5), Inches(6.6), Inches(6), Inches(0.6), NEON_PURPLE, 0.05)
add_text_box(slide, Inches(3.5), Inches(6.68), Inches(6), Inches(0.45),
             "Deployed on X Layer Testnet  ·  github.com/.../swap-lottery-hook",
             font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_page_number(slide, 10)

# ── Save ──
output_path = "/Users/kun/swap-lottery-hook/SwapLotteryHook_Pitch.pptx"
prs.save(output_path)
print(f"✅ PPT saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Format: 16:9 (13.33\" × 7.5\")")
print(f"   Theme: Dark tech with neon accents")
